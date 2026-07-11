"""
PowerPoint Updater Module
Updates presentation slides with KPI data.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import logging

logger = logging.getLogger(__name__)


class PowerPointUpdater:
    """Updates PowerPoint presentations."""

    def __init__(self, template_path=None):
        """
        Initialize PowerPoint updater.
        
        Args:
            template_path: Path to presentation template
        """
        self.template_path = template_path
        self.presentation = None

    def load_template(self, template_path=None):
        """
        Load presentation template.
        
        Args:
            template_path: Path to .pptx file
            
        Returns:
            Presentation object or None if error
        """
        try:
            path = template_path or self.template_path
            if not path:
                logger.error("No template path provided")
                return None

            path = Path(path)
            if not path.exists():
                logger.error(f"Template not found: {path}")
                return None

            self.presentation = Presentation(str(path))
            logger.info(f"Loaded template: {path.name}")
            return self.presentation

        except Exception as e:
            logger.error(f"Error loading template: {str(e)}")
            return None

    def create_blank(self):
        """Create blank presentation."""
        self.presentation = Presentation()
        logger.info("Created blank presentation")
        return self.presentation

    def get_slide_count(self):
        """Get number of slides."""
        if not self.presentation:
            return 0
        return len(self.presentation.slides)

    def get_slide(self, index):
        """Get slide by index."""
        if not self.presentation or index >= len(self.presentation.slides):
            return None
        return self.presentation.slides[index]

    def add_slide_with_title(self, title, subtitle=None):
        """
        Add slide with title and optional subtitle.
        
        Args:
            title: Slide title
            subtitle: Optional subtitle
            
        Returns:
            Slide object
        """
        try:
            slide_layout = self.presentation.slide_layouts[0]
            slide = self.presentation.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            title_shape.text = title

            if subtitle and len(slide.placeholders) > 1:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = subtitle

            logger.info(f"Added slide: {title}")
            return slide

        except Exception as e:
            logger.error(f"Error adding slide: {str(e)}")
            return None

    def update_text_in_slide(self, slide_index, placeholder_index, text):
        """
        Update text in slide placeholder.
        
        Args:
            slide_index: Slide index
            placeholder_index: Placeholder index
            text: New text
            
        Returns:
            True if successful
        """
        try:
            slide = self.get_slide(slide_index)
            if not slide or placeholder_index >= len(slide.placeholders):
                logger.error(f"Invalid slide or placeholder index")
                return False

            placeholder = slide.placeholders[placeholder_index]
            placeholder.text = str(text)
            logger.info(f"Updated text in slide {slide_index}, placeholder {placeholder_index}")
            return True

        except Exception as e:
            logger.error(f"Error updating text: {str(e)}")
            return False

    def add_text_box(self, slide_index, left, top, width, height, text):
        """
        Add text box to slide.
        
        Args:
            slide_index: Slide index
            left, top, width, height: Position and size in inches
            text: Text content
            
        Returns:
            TextFrame object or None
        """
        try:
            slide = self.get_slide(slide_index)
            if not slide:
                return None

            text_box = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            text_frame = text_box.text_frame
            text_frame.text = str(text)
            text_frame.word_wrap = True

            logger.info(f"Added text box to slide {slide_index}")
            return text_frame

        except Exception as e:
            logger.error(f"Error adding text box: {str(e)}")
            return None

    def add_table(self, slide_index, left, top, width, height, rows, cols, data=None):
        """
        Add table to slide.
        
        Args:
            slide_index: Slide index
            left, top, width, height: Position and size in inches
            rows: Number of rows
            cols: Number of columns
            data: Optional 2D list of cell data
            
        Returns:
            Table object or None
        """
        try:
            slide = self.get_slide(slide_index)
            if not slide:
                return None

            table_shape = slide.shapes.add_table(
                rows, cols,
                Inches(left), Inches(top),
                Inches(width), Inches(height)
            ).table

            if data:
                for i, row in enumerate(data):
                    for j, cell_value in enumerate(row):
                        if i < rows and j < cols:
                            table_shape.cell(i, j).text = str(cell_value)

            logger.info(f"Added table to slide {slide_index}")
            return table_shape

        except Exception as e:
            logger.error(f"Error adding table: {str(e)}")
            return None

    def add_image(self, slide_index, image_path, left, top, width=None, height=None):
        """
        Add image to slide.
        
        Args:
            slide_index: Slide index
            image_path: Path to image file
            left, top: Position in inches
            width, height: Optional dimensions in inches
            
        Returns:
            True if successful
        """
        try:
            slide = self.get_slide(slide_index)
            if not slide or not Path(image_path).exists():
                logger.error(f"Invalid slide or image path")
                return False

            if width and height:
                slide.shapes.add_picture(
                    str(image_path), Inches(left), Inches(top),
                    width=Inches(width), height=Inches(height)
                )
            else:
                slide.shapes.add_picture(
                    str(image_path), Inches(left), Inches(top)
                )

            logger.info(f"Added image to slide {slide_index}")
            return True

        except Exception as e:
            logger.error(f"Error adding image: {str(e)}")
            return False

    def format_text(self, text_frame, font_size=12, bold=False, color=None):
        """
        Format text in text frame.
        
        Args:
            text_frame: TextFrame object
            font_size: Font size in points
            bold: Bold flag
            color: RGB color tuple (r, g, b)
        """
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    run.font.bold = bold
                    if color:
                        run.font.color.rgb = RGBColor(*color)
        except Exception as e:
            logger.error(f"Error formatting text: {str(e)}")

    def save_presentation(self, output_path):
        """
        Save presentation to file.
        
        Args:
            output_path: Output file path
            
        Returns:
            True if successful
        """
        try:
            if not self.presentation:
                logger.error("No presentation loaded")
                return False

            output_path = Path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)

            self.presentation.save(str(output_path))
            logger.info(f"Saved presentation: {output_path}")
            return True

        except Exception as e:
            logger.error(f"Error saving presentation: {str(e)}")
            return False

    def update_kpi_slides(self, kpis, slide_start_index=1):
        """
        Update slides with KPI data.
        
        Args:
            kpis: Dictionary of KPIs
            slide_start_index: Starting slide index
            
        Returns:
            Number of slides updated
        """
        try:
            updated = 0
            for idx, (kpi_name, kpi_data) in enumerate(kpis.items()):
                slide_index = slide_start_index + idx
                if slide_index < len(self.presentation.slides):
                    value = kpi_data.get('value', 'N/A')
                    self.update_text_in_slide(slide_index, 1, f"{kpi_name}: {value}")
                    updated += 1

            logger.info(f"Updated {updated} KPI slides")
            return updated

        except Exception as e:
            logger.error(f"Error updating KPI slides: {str(e)}")
            return 0
